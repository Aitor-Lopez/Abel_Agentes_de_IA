import os
import uuid
import threading
import json
import markdown
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import HTMLResponse, FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
import subprocess
import tempfile
from evidence_reader import leer_evidencias, formatear_evidencias_para_prompt

from crewai import LLM, Agent, Task, Crew, Process
from crewai_tools import DirectoryReadTool, FileReadTool

import pdfplumber
from crewai.tools import BaseTool

class PDFReadTool(BaseTool):
    name: str = "PDFReadTool"
    description: str = "Lee archivos PDF y devuelve su contenido como texto. Úsalo para archivos .pdf"

    def _run(self, file_path: str) -> str:
        try:
            with pdfplumber.open(file_path) as pdf:
                return "\n".join(p.extract_text() or "" for p in pdf.pages)
        except Exception as e:
            return f"Error al leer el PDF: {e}"
        
load_dotenv()

app = FastAPI(title="AuditAI Server")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

EVIDENCIAS_DIR = Path("/app/evidencias")
DB_DIR = Path("/app/db")
EVIDENCIAS_DIR.mkdir(parents=True, exist_ok=True)
DB_DIR.mkdir(parents=True, exist_ok=True)

# Estado global de los jobs en memoria
jobs: dict[str, dict] = {}

# ──────────────────────────────────────────────────────────────
#  CONSTRUCCIÓN DEL CREW
# ──────────────────────────────────────────────────────────────

def build_crew() -> tuple[Crew, dict]:
    """
    Construye el crew completo y devuelve (crew, referencias_tareas).
    Las referencias a tareas permiten extraer sus outputs individuales.
    """
    llm = LLM(
        model=os.getenv("MODEL_NAME"),
        base_url=os.getenv("OLLAMA_BASE_URL"),
        temperature=0.1,
    )

    # 1. LEER TODAS LAS EVIDENCIAS CON PYTHON DIRECTAMENTE
    evidencias_dict = leer_evidencias(EVIDENCIAS_DIR)
    texto_evidencias = formatear_evidencias_para_prompt(evidencias_dict)

    # ── Agentes Auditores Especializados ──
    normativas = [
        {
            "n": "ISO 27001:2022",
            "g": "Analizar ÚNICAMENTE el contenido extraído de los archivos. E Identificar incumplimientos en los 93 controles del Anexo A de ISO 27001:2022, evaluando gestión de activos, control de accesos, criptografía, seguridad física, gestión de incidentes y continuidad de negocio.",
            "backstory": (
                "Eres un auditor extremadamente riguroso. TIENES PROHIBIDO inventar contenido. Si no usas las herramientas de lectura, no puedes saber qué hay en los archivos."
                "Eres un Lead Auditor certificado ISO/IEC 27001 con más de 12 años de experiencia en auditorías de SGSI "
                "para organizaciones del sector financiero, sanitario y de infraestructuras críticas. Has liderado más de "
                "80 proyectos de certificación en Europa y Latinoamérica. Tu metodología combina análisis de brechas "
                "técnicas con evaluación de madurez de procesos. Eres meticuloso, citas evidencias concretas y siempre "
                "relacionas los hallazgos con los controles específicos de la norma (A.5 a A.8)."
            ),
        },
        {
            "n": "ENS (Esquema Nacional de Seguridad)",
            "g": "Analizar ÚNICAMENTE el contenido extraído de los archivos. Y Auditar el cumplimiento del Real Decreto 311/2022 (ENS), evaluando las medidas de protección en categorías BÁSICA, MEDIA y ALTA.",
            "backstory": (
                "Eres un auditor extremadamente riguroso. TIENES PROHIBIDO inventar contenido. Si no usas las herramientas de lectura, no puedes saber qué hay en los archivos."
                "Eres un auditor acreditado ENS por el CCN-CERT con experiencia en administraciones públicas y "
                "proveedores de servicios digitales para el sector público español. Conoces en profundidad la Guía "
                "CCN-STIC 808. Has participado en más de 40 proyectos de adecuación ENS."
            ),
        },
        {
            "n": "GDPR / RGPD",
            "g": "Analizar ÚNICAMENTE el contenido extraído de los archivos. Y Verificar el cumplimiento del Reglamento (UE) 2016/679, analizando bases legales de tratamiento, derechos ARCO+, registro de actividades, EIPD y transferencias internacionales.",
            "backstory": (
                "Eres un auditor extremadamente riguroso. TIENES PROHIBIDO inventar contenido. Si no usas las herramientas de lectura, no puedes saber qué hay en los archivos."
                "Eres un Data Protection Officer (DPO) certificado con especialización en derecho digital europeo. "
                "Has asesorado a más de 60 organizaciones en su adecuación al RGPD y has gestionado notificaciones "
                "de brechas ante la AEPD. Dominas el análisis de flujos de datos y la elaboración de EIPD."
            ),
        },
        {
            "n": "NIS2 (Directiva UE 2022/2555)",
            "g": "Analizar ÚNICAMENTE el contenido extraído de los archivos. Y Evaluar la resiliencia operacional y obligaciones de gestión de riesgos, notificación de incidentes y seguridad de la cadena de suministro conforme a NIS2.",
            "backstory": (
                "Eres un auditor extremadamente riguroso. TIENES PROHIBIDO inventar contenido. Si no usas las herramientas de lectura, no puedes saber qué hay en los archivos."
                "Eres un experto en ciberseguridad de infraestructuras críticas con certificación CISSP. "
                "Has participado en ejercicios de resiliencia a nivel europeo coordinados por ENISA y conoces "
                "en detalle los requisitos de notificación de incidentes de NIS2."
            ),
        },
        {
            "n": "PCI-DSS v4.0",
            "g": "Analizar ÚNICAMENTE el contenido extraído de los archivos. Y Auditar los 12 requisitos de PCI-DSS v4.0 para proteger datos de tarjetas de pago, evaluando segmentación de red, cifrado, gestión de vulnerabilidades y controles de acceso.",
            "backstory": (
                "Eres un auditor extremadamente riguroso. TIENES PROHIBIDO inventar contenido. Si no usas las herramientas de lectura, no puedes saber qué hay en los archivos."
                "Eres un Qualified Security Assessor (QSA) certificado PCI-DSS con más de 10 años evaluando "
                "entornos CDE. Has realizado assessments para bancos y procesadores de pago de nivel 1. "
                "Conoces a fondo los cambios en la versión 4.0."
            ),
        },
    ]

    especialistas = [
        Agent(
            role=f"Auditor Senior {item['n']}",
            goal=item["g"],
            backstory=item["backstory"],
            llm=llm,
            verbose=True,
            max_iter=30,
            memory=False,
            allow_delegation=False,
            reasoning=False
        )
        for item in normativas
    ]

    tareas_audit = []

    for agent in especialistas:
        norma = agent.role.replace("Auditor Senior ", "")
        tareas_audit.append(
            Task(
                # 3. SIMPLIFICAR EL PROMPT E INYECTAR EL TEXTO
                description=f"""
                Tu objetivo es auditar la normativa {norma} basándote ÚNICAMENTE en las siguientes evidencias que ya han sido leídas y extraídas para ti:

                {texto_evidencias}

                REGLAS OBLIGATORIAS:
                1. Analiza exhaustivamente el texto de las evidencias proporcionadas arriba.
                2. Identifica controles cumplidos e incumplimientos.
                3. Marca “SIN EVIDENCIA” si no hay soporte documental para un control.
                4. Cada hallazgo debe incluir:
                   - control evaluado
                   - criticidad (Crítico/Alto/Medio/Bajo)
                   - impacto
                   - recomendación
                   - cita textual del archivo fuente

                PROHIBIDO:
                - Inventar datos, controles o evidencias que no estén en el texto superior.
                """,
                expected_output=(
                    "Un informe detallado de hallazgos basado puramente en las evidencias proporcionadas. "
                    "Debe incluir el control evaluado, estado (Cumple/Incumple/Sin evidencia), criticidad, "
                    "impacto, recomendación y la cita textual exacta del archivo que lo soporta."
                ),
                agent=agent,
                context=[],
            )
        )

    # ── Analista de Riesgos ──
    clasificador = Agent(
        role="Analista Senior de Riesgos Corporativos",
        goal=(
            "Consolidar todos los hallazgos de auditoría de las 5 normativas, eliminar duplicidades, "
            "correlacionar riesgos transversales y producir un registro unificado priorizado con scoring."
        ),
        backstory=(
            "Eres un Risk Manager certificado CRISC con experiencia en marcos ISO 31000, FAIR y NIST RMF. "
            "Has diseñado cuadros de mando de riesgos para consejos de administración. "
            "Tu valor diferencial es detectar riesgos transversales que afectan a múltiples normativas."
        ),
        llm=llm,
        verbose=True,
        max_iter=4,
        memory=False,
    )
    tarea_clasificacion = Task(
        description="""
Consolida los hallazgos de las 5 normativas. Elimina duplicidades, identifica riesgos transversales,
asigna scoring (criticidad × probabilidad × impacto, escala 1-25) y crea el registro unificado
ordenado por score descendente. Incluye métricas globales de exposición y nivel de madurez (1-5).
        """,
        expected_output=(
            "Registro unificado de riesgos con: métricas globales, top 10 riesgos críticos con score, "
            "registro completo, riesgos transversales y distribución por dominio."
        ),
        agent=clasificador,
        context=tareas_audit,
    )

    # ── Agente CISO ──
    agente_ciso = Agent(
        role="Redactor de Informe Técnico para CISO",
        goal="Generar informe técnico detallado para el CISO con hallazgos técnicos, análisis de brechas, métricas y roadmap de remediación.",
        backstory=(
            "Eres un consultor senior de ciberseguridad especializado en redactar informes técnicos para CISOs. "
            "Aportas profundidad técnica: controles exactos que fallan, arquitecturas deficientes y herramientas "
            "concretas a implementar. Tu estilo es directo, con tablas y métricas precisas."
        ),
        llm=llm,
        verbose=True,
        max_iter=4,
        memory=False,
    )
    tarea_informe_ciso = Task(
        description="""
Redacta el INFORME TÉCNICO PARA CISO en Markdown. Incluye:
- Resumen ejecutivo técnico con semáforo por normativa (🔴🟡🟢) e indicadores numéricos
- Análisis de brechas por normativa con % de cumplimiento y referencias exactas a controles
- Tabla completa de hallazgos: ID | Hallazgo | Normativas | Criticidad | Score | Evidencia | Recomendación Técnica
- Riesgos transversales y su impacto amplificado
- Gaps arquitectónicos y controles compensatorios
- Roadmap de remediación: Acción | Responsable | Herramienta | Plazo | KPI | Normativas que remedia
  * Fase 1 (0-30d): acciones críticas y quick-wins
  * Fase 2 (30-90d): mejoras estructurales
  * Fase 3 (90-180d): madurez y automatización
- KPIs de seguimiento y estimación de esfuerzo/presupuesto
        """,
        expected_output="Informe técnico Markdown para CISO con profundidad técnica, tablas detalladas y roadmap accionable.",
        agent=agente_ciso,
        context=[tarea_clasificacion, *tareas_audit],
    )

    # ── Agente CEO ──
    agente_ceo = Agent(
        role="Redactor de Informe Ejecutivo para CEO y Dirección",
        goal="Generar informe ejecutivo para CEO y dirección no técnica, traduciendo riesgos en impacto de negocio y decisiones estratégicas.",
        backstory=(
            "Eres un consultor de riesgos corporativos que presenta ante consejos de administración. "
            "Traduces riesgos técnicos a impacto económico, sanciones regulatorias y riesgo reputacional. "
            "Evitas jerga técnica y siempre cuantificas el impacto en términos de negocio."
        ),
        llm=llm,
        verbose=True,
        max_iter=4,
        memory=False,
    )
    tarea_informe_ceo = Task(
        description="""
Redacta el INFORME EJECUTIVO PARA CEO Y DIRECCIÓN en Markdown. SIN tecnicismos. Incluye:
- Resumen para la dirección (máximo 1 página): estado actual, semáforo normativo, decisión más urgente
- Exposición al riesgo regulatorio por normativa con cuantificación:
  * GDPR: hasta 20M€ o 4% facturación global
  * ENS: pérdida de contratos con AAPP
  * NIS2: multas hasta 10M€ o 2% facturación
  * PCI-DSS: revocación de capacidad de procesar pagos
  * ISO 27001: impacto en certificación y contratos
- Impacto potencial en el negocio: continuidad, reputación, clientes
- Los 5 riesgos más críticos explicados en lenguaje de negocio
- Tabla: Área de Mejora | Inversión Estimada | Riesgo Mitigado | ROI de Seguridad
- Decisiones estratégicas requeridas por la dirección
- Hoja de ruta ejecutiva a 6 meses (timeline simplificado)
- Top 3 acciones inmediatas (próximos 30 días)
        """,
        expected_output="Informe ejecutivo Markdown para CEO/Dirección sin tecnicismos, con foco en impacto de negocio y decisiones estratégicas.",
        agent=agente_ceo,
        context=[tarea_clasificacion],
    )

    # ── Agente General/DPO ──
    agente_general = Agent(
        role="Redactor de Informe General de Cumplimiento Normativo",
        goal="Generar informe formal de cumplimiento para DPO, compliance officers y auditores externos con matrices de controles y plan de acción.",
        backstory=(
            "Eres un especialista en compliance y gobierno corporativo. Has preparado informes para la AEPD, "
            "el CCN y auditores ISO. Aportas trazabilidad completa: cada hallazgo enlazado con el artículo "
            "exacto, la evidencia y el estado de remediación. Estilo formal y riguroso."
        ),
        llm=llm,
        verbose=True,
        max_iter=4,
        memory=False,
    )
    tarea_informe_general = Task(
        description="""
Redacta el INFORME GENERAL DE CUMPLIMIENTO NORMATIVO en Markdown. Incluye:
- Control del documento (versión, fecha, clasificación, elaborado/revisado/aprobado por)
- Alcance y metodología de la auditoría, limitaciones del análisis
- Estado de cumplimiento por normativa con tablas:
  | Control/Artículo | Descripción | Estado (C/NC/NE/NA) | Evidencia | Observaciones |
  (C=Cumple, NC=No Cumple, NE=No Evidencia, NA=No Aplica)
  Secciones para: ISO 27001, ENS, GDPR, NIS2, PCI-DSS
- Registro de No Conformidades: ID-NC | Normativa | Requisito | Criticidad | Fecha | Responsable | Estado
- Plan de Acción Correctivo (PAC): acción | responsable | fecha límite | indicador de cierre
- Análisis de riesgos para el Registro de Tratamientos GDPR (si hay datos personales en evidencias)
- Cadena de custodia de evidencias analizadas
- Declaraciones de aplicabilidad (controles ISO excluidos y justificación)
        """,
        expected_output="Informe formal Markdown para DPO/Compliance con matrices de controles, registro de no conformidades y PAC.",
        agent=agente_general,
        context=[tarea_clasificacion, *tareas_audit],
    )

    crew = Crew(
        agents=[*especialistas, clasificador, agente_ciso, agente_ceo, agente_general],
        tasks=[*tareas_audit, tarea_clasificacion, tarea_informe_ciso, tarea_informe_ceo, tarea_informe_general],
        process=Process.sequential,
        verbose=True,
    )

    task_refs = {
        "audit": tareas_audit,
        "clasificacion": tarea_clasificacion,
        "ciso": tarea_informe_ciso,
        "ceo": tarea_informe_ceo,
        "general": tarea_informe_general,
    }

    return crew, task_refs


def build_risk_ppt_crew(md_clasificacion: str, md_ciso: str) -> tuple:
    """
    Crew independiente para generar el contenido del PPT de análisis de riesgos.
    Recibe los MDs ya generados y produce un JSON estructurado para el PPTX.
    """
    llm = LLM(
        model=os.getenv("MODEL_NAME"),
        base_url=os.getenv("OLLAMA_BASE_URL"),
        temperature=0.15,
    )

    agente_ppt = Agent(
        role="Presentador Ejecutivo de Riesgos de Ciberseguridad",
        goal=(
            "Transformar los hallazgos de auditoría en una presentación PowerPoint de alto impacto "
            "orientada a vender la urgencia de los riesgos a la dirección. "
            "El objetivo es que los directivos entiendan el peligro real y aprueben inversiones."
        ),
        backstory=(
            "Eres un consultor senior de ciberseguridad con 15 años presentando ante consejos de administración. "
            "Dominas el arte de convertir datos técnicos en narrativas de riesgo que impactan. "
            "Sabes que los directivos responden a cifras concretas, consecuencias legales y comparaciones "
            "con casos reales del sector. Tu presentaciones siempre consiguen el presupuesto solicitado."
        ),
        llm=llm,
        verbose=True,
        max_iter=3,
        memory=False,
        allow_delegation=False,
        reasoning=False,
    )

    tarea_ppt = Task(
        description=f"""
Analiza los siguientes informes de auditoría y genera el contenido para una presentación PowerPoint
de ANÁLISIS DE RIESGOS orientada a vender la urgencia a la dirección.

=== INFORME DE CLASIFICACIÓN DE RIESGOS ===
{md_clasificacion[:8000]}

=== INFORME TÉCNICO CISO (primeras secciones) ===
{md_ciso[:6000]}

Genera un JSON EXACTAMENTE con esta estructura (sin texto adicional antes o después, solo el JSON):
{{
  "titulo": "Análisis de Riesgos de Ciberseguridad",
  "subtitulo": "Informe de Exposición y Riesgos Críticos",
  "fecha": "Abril 2026",
  "empresa": "AuditAI Multi-Agente",
  "resumen_ejecutivo": {{
    "nivel_riesgo": "CRÍTICO|ALTO|MEDIO|BAJO",
    "score_global": <número 1-100>,
    "hallazgos_criticos": <número>,
    "hallazgos_altos": <número>,
    "hallazgos_medios": <número>,
    "frase_impacto": "<frase contundente de máx 15 palabras sobre el riesgo>"
  }},
  "top_riesgos": [
    {{
      "id": "R1",
      "titulo": "<título corto del riesgo>",
      "descripcion": "<descripción en 2-3 frases orientada a impacto de negocio>",
      "impacto_economico": "<cuantificación en euros o % de facturación>",
      "probabilidad": "Alta|Media|Baja",
      "criticidad": "Crítico|Alto|Medio|Bajo",
      "normativas": ["ISO 27001", "GDPR"],
      "consecuencia": "<consecuencia directa si no se actúa>"
    }}
  ],
  "exposicion_normativa": [
    {{
      "normativa": "GDPR",
      "nivel": "CRÍTICO|ALTO|MEDIO|BAJO",
      "multa_max": "<multa máxima aplicable>",
      "incumplimientos": <número>,
      "descripcion": "<descripción breve del riesgo normativo>"
    }}
  ],
  "metricas_clave": [
    {{"metrica": "<nombre>", "valor": "<valor>", "contexto": "<qué significa>"}},
    {{"metrica": "<nombre>", "valor": "<valor>", "contexto": "<qué significa>"}},
    {{"metrica": "<nombre>", "valor": "<valor>", "contexto": "<qué significa>"}},
    {{"metrica": "<nombre>", "valor": "<valor>", "contexto": "<qué significa>"}}
  ],
  "plan_accion": [
    {{
      "fase": "Inmediato (0-30 días)",
      "acciones": ["<acción 1>", "<acción 2>", "<acción 3>"],
      "inversion": "<rango en euros>",
      "riesgo_mitigado": "<qué riesgo se elimina>"
    }},
    {{
      "fase": "Corto plazo (30-90 días)",
      "acciones": ["<acción 1>", "<acción 2>", "<acción 3>"],
      "inversion": "<rango en euros>",
      "riesgo_mitigado": "<qué riesgo se elimina>"
    }},
    {{
      "fase": "Medio plazo (90-180 días)",
      "acciones": ["<acción 1>", "<acción 2>", "<acción 3>"],
      "inversion": "<rango en euros>",
      "riesgo_mitigado": "<qué riesgo se elimina>"
    }}
  ],
  "llamada_accion": "<mensaje final contundente de máx 20 palabras para cerrar la presentación>"
}}

IMPORTANTE: Usa ÚNICAMENTE datos de los informes proporcionados. El JSON debe ser válido y parseable.
Incluye exactamente 5 elementos en top_riesgos y 5 en exposicion_normativa (ISO 27001, ENS, GDPR, NIS2, PCI-DSS).
Incluye exactamente 4 métricas clave.
        """,
        expected_output="JSON válido con la estructura exacta especificada, sin texto adicional.",
        agent=agente_ppt,
    )

    crew = Crew(
        agents=[agente_ppt],
        tasks=[tarea_ppt],
        process=Process.sequential,
        verbose=True,
    )

    return crew, tarea_ppt


# ──────────────────────────────────────────────────────────────
#  EJECUCIÓN EN BACKGROUND
# ──────────────────────────────────────────────────────────────

def run_analysis(job_id: str):
    """Ejecuta el crew en un thread separado y actualiza el estado del job."""
    import time

    job = jobs[job_id]

    agentes_orden = [
        ("iso",     "Auditor Senior ISO 27001:2022"),
        ("ens",     "Auditor Senior ENS"),
        ("gdpr",    "Auditor Senior GDPR / RGPD"),
        ("nis2",    "Auditor Senior NIS2"),
        ("pci",     "Auditor Senior PCI-DSS v4.0"),
        ("risk",    "Analista Senior de Riesgos"),
        ("ciso",    "Redactor Informe CISO"),
        ("ceo",     "Redactor Informe CEO / Dirección"),
        ("general", "Redactor Informe Cumplimiento General"),
    ]

    def set_agent(idx, state, msg=""):
        job["agente_actual"]  = agentes_orden[idx][0]
        job["agente_nombre"]  = agentes_orden[idx][1]
        job["agente_estado"]  = state
        job["agente_idx"]     = idx
        job["mensaje"]        = msg

    try:
        archivos = [f for f in EVIDENCIAS_DIR.iterdir() if f.is_file()]
        if not archivos:
            job["estado"] = "error"
            job["error"]  = "No hay archivos en la carpeta de evidencias."
            return

        job["estado"]        = "running"
        job["total_agentes"] = len(agentes_orden)

        crew, task_refs = build_crew()

        resultado_holder = {"result": None, "error": None, "tasks": None}

        def crew_thread():
            try:
                resultado_holder["result"] = crew.kickoff()
                resultado_holder["tasks"]  = task_refs
            except Exception as e:
                resultado_holder["error"] = str(e)

        t = threading.Thread(target=crew_thread)
        t.start()

        # Tiempos estimados por agente (segundos)
        tiempos = [45, 45, 45, 40, 40, 30, 35, 25, 35]

        for i, (agent_id, agent_name) in enumerate(agentes_orden):
            set_agent(i, "running", f"Analizando con {agent_name}…")
            wait = tiempos[i] if i < len(tiempos) else 30
            for _ in range(wait * 10):
                if not t.is_alive():
                    break
                time.sleep(0.1)
            if not t.is_alive():
                for j in range(i, len(agentes_orden)):
                    set_agent(j, "done", "")
                break
            set_agent(i, "done", "")

        t.join()

        if resultado_holder["error"]:
            job["estado"] = "error"
            job["error"]  = resultado_holder["error"]
            return

        tasks = resultado_holder["tasks"]

        # Extraer outputs de cada tarea
        def safe_output(task):
            if task and task.output:
                return task.output.raw or str(task.output)
            return ""

        md_ciso    = safe_output(tasks["ciso"])
        md_ceo     = safe_output(tasks["ceo"])
        md_general = safe_output(tasks["general"])
        md_risk    = safe_output(tasks["clasificacion"])

        # Si alguno está vacío, usar el resultado final del crew
        result_str = str(resultado_holder["result"])
        if not md_ciso:    md_ciso    = result_str
        if not md_ceo:     md_ceo     = result_str
        if not md_general: md_general = result_str

        # Guardar Markdown en disco
        informes_md = {
            "ciso":    md_ciso,
            "ceo":     md_ceo,
            "general": md_general,
            "clasificacion": md_risk,
        }
        for nombre, contenido in informes_md.items():
            if contenido:
                path = DB_DIR / f"informe_{nombre}.md"
                path.write_text(contenido, encoding="utf-8")

        # Convertir a HTML
        html_ciso    = markdown_to_html(md_ciso,    "CISO — Técnico",           "ciso")
        html_ceo     = markdown_to_html(md_ceo,     "CEO / Dirección — Ejecutivo", "ceo")
        html_general = markdown_to_html(md_general, "Cumplimiento General — DPO/Compliance", "general")

        for nombre, html in [("ciso", html_ciso), ("ceo", html_ceo), ("general", html_general)]:
            path = DB_DIR / f"informe_{nombre}.html"
            path.write_text(html, encoding="utf-8")

        job["estado"]       = "done"
        job["fecha"]        = datetime.now().isoformat()
        job["informes_md"]  = informes_md
        job["informes_html"] = {
            "ciso":    html_ciso,
            "ceo":     html_ceo,
            "general": html_general,
        }

    except Exception as e:
        job["estado"] = "error"
        job["error"]  = str(e)


# ──────────────────────────────────────────────────────────────
#  HTML RENDERER
# ──────────────────────────────────────────────────────────────

INFORME_LABELS = {
    "ciso":    ("Informe Técnico — CISO",                  "#5b9fff", "👨‍💻"),
    "ceo":     ("Informe Ejecutivo — CEO / Dirección",     "#00e5a0", "📊"),
    "general": ("Informe de Cumplimiento — DPO/Compliance","#a78bfa", "📋"),
    "riesgos": ("Análisis de Riesgos — Presentación Ejecutiva", "#ff4757", "⚠️"),
}

def markdown_to_html(md_text: str, subtitulo: str, tipo: str) -> str:
    body = markdown.markdown(
        md_text,
        extensions=["tables", "fenced_code", "nl2br", "toc"],
    )

    # Badges de criticidad
    for palabra, cls in [("Crítico", "critico"), ("Alto", "alto"), ("Medio", "medio"), ("Bajo", "bajo")]:
        body = body.replace(f">{palabra}<",        f'><span class="badge {cls}">{palabra}</span><')
        body = body.replace(f"**{palabra}**",      f'<span class="badge {cls}">{palabra}</span>')

    # Badges de estado de cumplimiento
    for palabra, cls in [("No Cumple", "nc"), ("Cumple", "c"), ("No Evidencia", "ne"), ("No Aplica", "na")]:
        body = body.replace(f">{palabra}<",   f'><span class="badge {cls}">{palabra}</span><')

    label, color, emoji = INFORME_LABELS.get(tipo, (subtitulo, "#00e5a0", "📄"))

    return f"""<!DOCTYPE html>
<html lang="es">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{label} — AuditAI</title>
<link href="https://fonts.googleapis.com/css2?family=Syne:wght@700;800&family=Inter:wght@400;500&family=IBM+Plex+Mono:wght@400;500&display=swap" rel="stylesheet">
<style>
  :root {{
    --bg: #0d1117; --bg2: #161b22; --bg3: #21262d;
    --accent: {color}; --danger: #ff4757; --warn: #ffa502;
    --info: #5b9fff; --purple: #a78bfa;
    --text: #e8eaf0; --text2: #8892a4; --text3: #4a5568;
    --border: #21262d;
  }}
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  html {{ overflow-x: hidden; }}
  body {{
    background: var(--bg); color: var(--text);
    font-family: 'Inter', sans-serif; font-size: 15px;
    line-height: 1.8; padding: 0;
    overflow-x: hidden;
  }}
  .page-wrap {{ max-width: 980px; margin: 0 auto; padding: 60px 40px 120px; }}

  h1 {{ font-family: 'Syne', sans-serif; font-size: 38px; font-weight: 800;
        color: var(--accent); margin-bottom: 8px; line-height: 1.1; letter-spacing: -1px; }}
  h2 {{ font-family: 'Syne', sans-serif; font-size: 21px; font-weight: 700;
        color: #fff; margin: 52px 0 16px; padding-bottom: 10px;
        border-bottom: 2px solid var(--border); }}
  h3 {{ font-family: 'Syne', sans-serif; font-size: 16px; font-weight: 700;
        color: var(--accent); margin: 32px 0 10px; }}
  h4 {{ font-size: 14px; font-weight: 600; color: var(--text2); margin: 20px 0 8px; text-transform: uppercase; letter-spacing: 0.5px; }}
  p {{ margin-bottom: 14px; color: var(--text); }}
  a {{ color: var(--accent); }}
  strong {{ color: #fff; font-weight: 500; }}
  em {{ color: var(--text2); font-style: italic; }}
  blockquote {{
    border-left: 3px solid var(--accent); margin: 24px 0;
    padding: 14px 20px; background: rgba(255,255,255,0.03);
    border-radius: 0 8px 8px 0; color: var(--text2);
  }}
  code {{
    background: var(--bg2); color: var(--accent);
    padding: 2px 7px; border-radius: 4px;
    font-family: 'IBM Plex Mono', monospace; font-size: 13px;
  }}
  pre {{ background: var(--bg2); border: 1px solid var(--border);
         border-radius: 10px; padding: 20px; overflow-x: auto; margin: 20px 0; }}
  pre code {{ background: none; padding: 0; color: var(--text); }}
  ul, ol {{ padding-left: 24px; margin-bottom: 14px; }}
  li {{ margin-bottom: 6px; }}
  hr {{ border: none; border-top: 1px solid var(--border); margin: 48px 0; }}

  /* TABLES */
  .table-wrap {{ overflow-x: auto; margin: 24px 0; border-radius: 12px;
                  border: 1px solid var(--border); max-width: 100%; }}
  table {{ width: 100%; border-collapse: collapse; font-size: 13px; table-layout: auto; }}
  thead tr {{ background: var(--bg2); }}
  th {{ padding: 10px 12px; text-align: left; color: var(--text2);
        font-size: 10px; font-weight: 500; text-transform: uppercase;
        letter-spacing: 1px; white-space: nowrap; }}
  td {{ padding: 10px 12px; border-top: 1px solid var(--border); vertical-align: top;
       overflow-wrap: break-word; word-break: normal; hyphens: auto; min-width: 80px; }}
  td:first-child {{ min-width: 120px; max-width: 200px; }}
  tbody tr:hover {{ background: rgba(255,255,255,0.02); }}

  /* BADGES — criticidad */
  .badge {{ display: inline-block; padding: 2px 10px; border-radius: 20px;
             font-size: 11px; font-weight: 600; letter-spacing: 0.5px; white-space: nowrap; }}
  .critico {{ background: rgba(255,71,87,0.15); color: #ff6b7a; border: 1px solid rgba(255,71,87,0.3); }}
  .alto    {{ background: rgba(255,165,2,0.15);  color: #ffc04d; border: 1px solid rgba(255,165,2,0.3); }}
  .medio   {{ background: rgba(255,215,0,0.12);  color: #ffd700; border: 1px solid rgba(255,215,0,0.25); }}
  .bajo    {{ background: rgba(0,229,160,0.12);   color: #00e5a0; border: 1px solid rgba(0,229,160,0.25); }}

  /* BADGES — estado cumplimiento */
  .c  {{ background: rgba(0,229,160,0.12); color: #00e5a0;  border: 1px solid rgba(0,229,160,0.3); }}
  .nc {{ background: rgba(255,71,87,0.15); color: #ff6b7a;  border: 1px solid rgba(255,71,87,0.3); }}
  .ne {{ background: rgba(255,165,2,0.12); color: #ffc04d;  border: 1px solid rgba(255,165,2,0.3); }}
  .na {{ background: rgba(255,255,255,0.06); color: #8892a4; border: 1px solid rgba(255,255,255,0.1); }}

  /* REPORT HEADER */
  .report-header {{
    background: var(--bg2); border: 1px solid var(--border);
    border-radius: 16px; padding: 36px 40px; margin-bottom: 52px;
    position: relative; overflow: hidden;
  }}
  .report-header::before {{
    content: '{emoji}';
    position: absolute; right: 36px; top: 50%; transform: translateY(-50%);
    font-size: 64px; opacity: 0.15;
  }}
  .report-type {{
    font-family: 'IBM Plex Mono', monospace; font-size: 11px;
    color: var(--accent); letter-spacing: 2px; text-transform: uppercase;
    margin-bottom: 10px;
  }}
  .report-meta {{
    display: flex; gap: 24px; flex-wrap: wrap; margin-top: 16px;
  }}
  .report-meta span {{
    font-size: 12px; color: var(--text3);
    font-family: 'IBM Plex Mono', monospace; letter-spacing: 0.3px;
  }}
  .report-meta span b {{ color: var(--text2); font-weight: 500; }}

  @media (max-width: 640px) {{
    .page-wrap {{ padding: 32px 20px 60px; }}
    h1 {{ font-size: 26px; }}
    h2 {{ font-size: 18px; }}
  }}
</style>
</head>
<body>
<div class="page-wrap">
  <div class="report-header">
    <div class="report-type">AuditAI — Informe de Seguridad</div>
    <h1>{label}</h1>
    <div class="report-meta">
      <span><b>Sistema:</b> AuditAI Multi-Agente (CrewAI + Ollama)</span>
      <span><b>Normativas:</b> ISO 27001 · ENS · GDPR · NIS2 · PCI-DSS</span>
      <span><b>Destinatario:</b> {subtitulo}</span>
    </div>
  </div>
  {body}
</div>
<script>
  document.querySelectorAll('table').forEach(t => {{
    if (!t.parentElement.classList.contains('table-wrap')) {{
      const w = document.createElement('div');
      w.className = 'table-wrap';
      t.parentNode.insertBefore(w, t);
      w.appendChild(t);
    }}
  }});
</script>
</body>
</html>"""


# ──────────────────────────────────────────────────────────────
#  ENDPOINTS
# ──────────────────────────────────────────────────────────────

@app.get("/", response_class=HTMLResponse)
async def index():
    html_path = Path("/app/static/index.html")
    if html_path.exists():
        return HTMLResponse(html_path.read_text(encoding="utf-8"))
    return HTMLResponse("<h1>AuditAI</h1><p>Falta /app/static/index.html</p>")


@app.get("/api/evidencias")
async def listar_evidencias():
    archivos = []
    for f in EVIDENCIAS_DIR.iterdir():
        if f.is_file():
            archivos.append({
                "nombre":     f.name,
                "tamaño":     f.stat().st_size,
                "modificado": datetime.fromtimestamp(f.stat().st_mtime).isoformat(),
            })
    return {"archivos": sorted(archivos, key=lambda x: x["nombre"])}


@app.post("/api/evidencias/upload")
async def subir_evidencias(files: list[UploadFile] = File(...)):
    subidos = []
    for file in files:
        dest    = EVIDENCIAS_DIR / file.filename
        content = await file.read()
        dest.write_bytes(content)
        subidos.append({"nombre": file.filename, "tamaño": len(content)})
    return {"subidos": subidos}


@app.delete("/api/evidencias/{nombre}")
async def borrar_evidencia(nombre: str):
    path = EVIDENCIAS_DIR / nombre
    if not path.exists():
        raise HTTPException(404, "Archivo no encontrado")
    path.unlink()
    return {"ok": True}


@app.post("/api/analizar")
async def iniciar_analisis():
    archivos = [f for f in EVIDENCIAS_DIR.iterdir() if f.is_file()]
    if not archivos:
        raise HTTPException(400, "No hay archivos en la carpeta de evidencias.")

    job_id = str(uuid.uuid4())
    jobs[job_id] = {
        "estado":        "starting",
        "agente_actual": None,
        "agente_nombre": None,
        "agente_estado": None,
        "agente_idx":    -1,
        "total_agentes": 9,
        "mensaje":       "Preparando agentes…",
        "error":         None,
        "informes_md":   None,
        "informes_html": None,
        "fecha":         None,
    }
    threading.Thread(target=run_analysis, args=(job_id,), daemon=True).start()
    return {"job_id": job_id}


@app.get("/api/estado/{job_id}")
async def estado_job(job_id: str):
    if job_id not in jobs:
        raise HTTPException(404, "Job no encontrado")
    job = jobs[job_id]
    return {
        "estado":        job["estado"],
        "agente_idx":    job["agente_idx"],
        "agente_nombre": job["agente_nombre"],
        "agente_estado": job["agente_estado"],
        "total_agentes": job["total_agentes"],
        "mensaje":       job["mensaje"],
        "error":         job.get("error"),
        "fecha":         job.get("fecha"),
        "informes_disponibles": (
            list(job["informes_html"].keys())
            if job.get("informes_html") else []
        ),
    }


# ── Rutas estáticas PRIMERO (antes que las dinámicas con {job_id}) ──
# FastAPI resuelve rutas en orden de registro. Si las rutas con {job_id}
# fueran primero, "ultimo" sería capturado como job_id y devolvería 404.

@app.get("/api/informe/ultimo/{tipo}/html", response_class=HTMLResponse)
async def ultimo_informe_html(tipo: str):
    """Último informe guardado en disco. tipo: ciso | ceo | general"""
    path = DB_DIR / f"informe_{tipo}.html"
    if not path.exists():
        raise HTTPException(404, f"No hay informe '{tipo}' generado todavía")
    return HTMLResponse(path.read_text(encoding="utf-8"))


@app.get("/api/informe/ultimo/{tipo}/descargar")
async def descargar_informe(tipo: str):
    """Descarga el informe HTML como archivo."""
    path = DB_DIR / f"informe_{tipo}.html"
    if not path.exists():
        raise HTTPException(404, f"No hay informe '{tipo}' generado todavía")
    fecha = datetime.now().strftime("%Y%m%d")
    return FileResponse(
        path, media_type="text/html",
        filename=f"auditoria_{tipo}_{fecha}.html",
    )


# ── Rutas dinámicas con {job_id} DESPUÉS ──

@app.get("/api/informe/{job_id}/{tipo}/html", response_class=HTMLResponse)
async def obtener_informe_html(job_id: str, tipo: str):
    """Devuelve un informe específico como HTML. tipo: ciso | ceo | general"""
    if job_id not in jobs:
        raise HTTPException(404, "Job no encontrado")
    job = jobs[job_id]
    if job["estado"] != "done":
        raise HTTPException(400, "El informe aún no está listo")
    if not job.get("informes_html") or tipo not in job["informes_html"]:
        raise HTTPException(404, f"Informe '{tipo}' no disponible")
    return HTMLResponse(job["informes_html"][tipo])


@app.get("/api/informe/{job_id}/{tipo}/md")
async def obtener_informe_md(job_id: str, tipo: str):
    """Devuelve un informe específico como Markdown. tipo: ciso | ceo | general | clasificacion"""
    if job_id not in jobs:
        raise HTTPException(404, "Job no encontrado")
    job = jobs[job_id]
    if job["estado"] != "done":
        raise HTTPException(400, "El informe aún no está listo")
    if not job.get("informes_md") or tipo not in job["informes_md"]:
        raise HTTPException(404, f"Informe '{tipo}' no disponible")
    return {"markdown": job["informes_md"][tipo]}


@app.get("/api/informes/disponibles")
async def informes_disponibles():
    """Lista qué informes están guardados en disco."""
    tipos = ["ciso", "ceo", "general", "clasificacion"]
    resultado = {}
    for tipo in tipos:
        md_path   = DB_DIR / f"informe_{tipo}.md"
        html_path = DB_DIR / f"informe_{tipo}.html"
        resultado[tipo] = {
            "md":   md_path.exists(),
            "html": html_path.exists(),
            "fecha": (
                datetime.fromtimestamp(html_path.stat().st_mtime).isoformat()
                if html_path.exists() else None
            ),
        }
    resultado["riesgos_ppt"] = {
        "pptx": (DB_DIR / "informe_riesgos.pptx").exists(),
        "json": (DB_DIR / "informe_riesgos.json").exists(),
        "html": (DB_DIR / "informe_riesgos.html").exists(),
        "fecha": (
            datetime.fromtimestamp((DB_DIR / "informe_riesgos.pptx").stat().st_mtime).isoformat()
            if (DB_DIR / "informe_riesgos.pptx").exists() else None
        ),
    }
    return resultado


# ──────────────────────────────────────────────────────────────
#  PPT DE RIESGOS — jobs separados
# ──────────────────────────────────────────────────────────────

ppt_jobs: dict[str, dict] = {}


def run_risk_ppt(job_id: str):
    """Ejecuta el agente PPT en background y genera el PPTX con python-pptx."""
    job = ppt_jobs[job_id]
    try:
        # Leer los MDs de disco
        md_clas_path = DB_DIR / "informe_clasificacion.md"
        md_ciso_path = DB_DIR / "informe_ciso.md"
        md_clas = md_clas_path.read_text(encoding="utf-8") if md_clas_path.exists() else ""
        md_ciso = md_ciso_path.read_text(encoding="utf-8") if md_ciso_path.exists() else ""

        if not md_clas and not md_ciso:
            job["estado"] = "error"
            job["error"] = "No hay informes generados. Ejecuta primero el análisis principal."
            return

        job["estado"] = "running"
        job["mensaje"] = "Agente analizando riesgos para la presentación…"

        crew, tarea = build_risk_ppt_crew(md_clas, md_ciso)
        resultado = crew.kickoff()

        raw = tarea.output.raw if tarea.output else str(resultado)

        # Extraer JSON limpio
        import re
        json_match = re.search(r'\{[\s\S]*\}', raw)
        if not json_match:
            raise ValueError("El agente no devolvió un JSON válido")

        data = json.loads(json_match.group())

        # Guardar JSON
        json_path = DB_DIR / "informe_riesgos.json"
        json_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

        # Generar PPTX con python-pptx
        pptx_path = DB_DIR / "informe_riesgos.pptx"
        generate_pptx_python(data, str(pptx_path))

        # Generar HTML de preview
        html_preview = risk_data_to_html(data)
        html_path = DB_DIR / "informe_riesgos.html"
        html_path.write_text(html_preview, encoding="utf-8")

        job["estado"] = "done"
        job["fecha"] = datetime.now().isoformat()
        job["data"] = data

    except Exception as e:
        job["estado"] = "error"
        job["error"] = str(e)


def _rgb(hex_color: str):
    """Convierte hex string (sin #) a RGBColor de pptx."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def generate_pptx_python(data: dict, output_path: str):
    """Genera el PPTX completo usando python-pptx (sin Node.js)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    # ── Paleta ──
    C = {
        "bg":       "0D1117",
        "bg2":      "161B22",
        "bg3":      "21262D",
        "red":      "FF4757",
        "orange":   "FFA502",
        "green":    "00E5A0",
        "blue":     "5B9FFF",
        "white":    "E8EAF0",
        "gray":     "8892A4",
        "darkgray": "4A5568",
    }
    CRIT_COLOR  = {"Crítico": C["red"], "Alto": C["orange"], "Medio": "FFD700", "Bajo": C["green"]}
    NIVEL_COLOR = {"CRÍTICO": C["red"], "ALTO":  C["orange"], "MEDIO": "FFD700", "BAJO": C["green"]}

    def rgb(k): return _rgb(C[k])
    def rgb_hex(h): return _rgb(h)

    prs = Presentation()
    prs.slide_width  = Inches(13.3)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # totalmente en blanco

    def add_slide():
        s = prs.slides.add_slide(blank_layout)
        # Fondo oscuro
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb("bg")
        return s

    def add_rect(slide, x, y, w, h, fill_hex, alpha=None):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb_hex(fill_hex)
        shape.line.fill.background()  # sin borde
        return shape

    def add_text(slide, text, x, y, w, h, font_size=12, bold=False,
                 color_hex=None, align="left", italic=False):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.text = str(text)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color_hex:
            run.font.color.rgb = rgb_hex(color_hex)
        else:
            run.font.color.rgb = rgb("white")
        return txBox

    resumen = data.get("resumen_ejecutivo", {})
    nivel = resumen.get("nivel_riesgo", "ALTO")
    nivel_color = NIVEL_COLOR.get(nivel, C["orange"])
    score = resumen.get("score_global", 0)

    # ══ SLIDE 1 — PORTADA ══
    s1 = add_slide()
    add_rect(s1, 0, 0, 0.5, 7.5, C["red"])
    add_rect(s1, 0.9, 1.0, 3.2, 0.65, nivel_color)
    add_text(s1, f"NIVEL DE RIESGO: {nivel}", 0.9, 1.0, 3.2, 0.65,
             font_size=14, bold=True, color_hex=C["bg"], align="center")
    add_text(s1, data.get("titulo", "Análisis de Riesgos"), 0.9, 1.85, 9, 1.3,
             font_size=36, bold=True, color_hex=C["white"])
    add_text(s1, data.get("subtitulo", ""), 0.9, 3.1, 9, 0.6,
             font_size=18, color_hex=nivel_color)
    add_rect(s1, 0.9, 3.9, 8.5, 0.8, C["bg2"])
    add_text(s1, f"⚠  {resumen.get('frase_impacto', '')}", 0.9, 3.9, 8.5, 0.8,
             font_size=14, color_hex=C["red"], italic=True, align="center")
    mets = [
        ("Hallazgos Críticos", str(resumen.get("hallazgos_criticos", 0)), C["red"]),
        ("Hallazgos Altos",    str(resumen.get("hallazgos_altos", 0)),    C["orange"]),
        ("Hallazgos Medios",   str(resumen.get("hallazgos_medios", 0)),   "FFD700"),
        ("Score Global",       f"{score}/100",                             nivel_color),
    ]
    for i, (label, val, col) in enumerate(mets):
        bx = 0.9 + i * 2.9
        add_rect(s1, bx, 4.9, 2.6, 1.3, C["bg2"])
        add_rect(s1, bx, 4.9, 2.6, 0.07, col)
        add_text(s1, val,   bx, 5.0, 2.6, 0.65, font_size=28, bold=True,  color_hex=col, align="center")
        add_text(s1, label, bx, 5.7, 2.6, 0.4,  font_size=9,  bold=False, color_hex=C["gray"], align="center")
    add_text(s1, f"{data.get('empresa','')}  ·  {data.get('fecha','')}",
             0.9, 7.0, 10, 0.35, font_size=9, color_hex=C["darkgray"])

    # ══ SLIDE 2 — TOP 5 RIESGOS ══
    s2 = add_slide()
    add_rect(s2, 0, 0, 13.3, 0.8, C["bg2"])
    add_text(s2, "TOP 5 RIESGOS CRÍTICOS", 0.5, 0, 8, 0.8, font_size=20, bold=True, color_hex=C["white"])
    add_text(s2, "Impacto directo sobre el negocio si no se actúa",
             0.5, 0, 12.3, 0.8, font_size=11, color_hex=C["gray"], align="right")
    for i, r in enumerate((data.get("top_riesgos") or [])[:5]):
        ry = 0.95 + i * 1.28
        col = CRIT_COLOR.get(r.get("criticidad", ""), C["orange"])
        add_rect(s2, 0.3, ry, 12.7, 1.15, C["bg2"])
        add_rect(s2, 0.3, ry, 0.08, 1.15, col)
        add_rect(s2, 0.55, ry + 0.25, 0.5, 0.5, col)
        add_text(s2, r.get("id", ""), 0.55, ry + 0.25, 0.5, 0.5,
                 font_size=11, bold=True, color_hex=C["bg"], align="center")
        add_text(s2, r.get("titulo", ""), 1.2, ry + 0.08, 5.5, 0.4,
                 font_size=13, bold=True, color_hex=C["white"])
        add_text(s2, r.get("descripcion", ""), 1.2, ry + 0.5, 5.5, 0.55,
                 font_size=9, color_hex=C["gray"])
        add_text(s2, f"💰 {r.get('impacto_economico','')}", 6.9, ry + 0.1, 3.5, 0.4,
                 font_size=10, bold=True, color_hex=col)
        add_text(s2, f"⚡ {r.get('consecuencia','')}", 6.9, ry + 0.55, 3.5, 0.45,
                 font_size=9, color_hex=C["gray"])
        add_rect(s2, 10.6, ry + 0.3, 1.5, 0.4, col)
        add_text(s2, r.get("criticidad", "").upper(), 10.6, ry + 0.3, 1.5, 0.4,
                 font_size=9, bold=True, color_hex=C["bg"], align="center")
        norms_txt = " · ".join(r.get("normativas") or [])
        add_text(s2, norms_txt, 10.55, ry + 0.75, 1.6, 0.3,
                 font_size=7, color_hex=C["darkgray"], align="center")

    # ══ SLIDE 3 — EXPOSICIÓN NORMATIVA ══
    s3 = add_slide()
    add_rect(s3, 0, 0, 13.3, 0.8, C["bg2"])
    add_text(s3, "EXPOSICIÓN NORMATIVA Y SANCIONES", 0.5, 0, 12, 0.8,
             font_size=20, bold=True, color_hex=C["white"])
    for i, n in enumerate((data.get("exposicion_normativa") or [])[:5]):
        col_i = i % 3
        row_i = i // 3
        cx = 0.4 + col_i * 4.3
        cy = 1.0 + row_i * 2.8
        cw, ch = 3.9, 2.5
        col = NIVEL_COLOR.get(n.get("nivel", ""), C["orange"])
        add_rect(s3, cx, cy, cw, ch, C["bg2"])
        add_rect(s3, cx, cy, cw, 0.07, col)
        add_text(s3, n.get("normativa", ""), cx + 0.15, cy + 0.12, cw - 0.6, 0.45,
                 font_size=15, bold=True, color_hex=C["white"])
        add_rect(s3, cx + cw - 0.9, cy + 0.12, 0.75, 0.35, col)
        add_text(s3, n.get("nivel", ""), cx + cw - 0.9, cy + 0.12, 0.75, 0.35,
                 font_size=7, bold=True, color_hex=C["bg"], align="center")
        add_text(s3, f"Multa máx: {n.get('multa_max','')}", cx + 0.15, cy + 0.65, cw - 0.3, 0.35,
                 font_size=10, bold=True, color_hex=col)
        add_text(s3, f"{n.get('incumplimientos',0)} incumplimientos detectados",
                 cx + 0.15, cy + 1.0, cw - 0.3, 0.3, font_size=9, color_hex=C["gray"])
        add_text(s3, n.get("descripcion", ""), cx + 0.15, cy + 1.35, cw - 0.3, 0.95,
                 font_size=8, color_hex=C["gray"])

    # ══ SLIDE 4 — MÉTRICAS CLAVE ══
    s4 = add_slide()
    add_rect(s4, 0, 0, 13.3, 0.8, C["bg2"])
    add_text(s4, "MÉTRICAS CLAVE DE EXPOSICIÓN", 0.5, 0, 12, 0.8,
             font_size=20, bold=True, color_hex=C["white"])
    m_colors = [C["red"], C["orange"], C["blue"], C["green"]]
    for i, m in enumerate((data.get("metricas_clave") or [])[:4]):
        col_i = i % 2
        row_i = i // 2
        mx = 0.5 + col_i * 6.4
        my = 1.0 + row_i * 2.8
        mw, mh = 6.0, 2.5
        mc = m_colors[i]
        add_rect(s4, mx, my, mw, mh, C["bg2"])
        add_rect(s4, mx, my, mw, 0.07, mc)
        add_text(s4, str(m.get("valor", "")), mx, my + 0.15, mw, 1.2,
                 font_size=44, bold=True, color_hex=mc, align="center")
        add_text(s4, m.get("metrica", ""), mx + 0.2, my + 1.4, mw - 0.4, 0.45,
                 font_size=14, bold=True, color_hex=C["white"], align="center")
        add_text(s4, m.get("contexto", ""), mx + 0.2, my + 1.88, mw - 0.4, 0.45,
                 font_size=10, color_hex=C["gray"], align="center")

    # ══ SLIDE 5 — PLAN DE ACCIÓN ══
    s5 = add_slide()
    add_rect(s5, 0, 0, 13.3, 0.8, C["bg2"])
    add_text(s5, "PLAN DE ACCIÓN RECOMENDADO", 0.5, 0, 12, 0.8,
             font_size=20, bold=True, color_hex=C["white"])
    fase_colors = [C["red"], C["orange"], C["blue"]]
    for i, f in enumerate((data.get("plan_accion") or [])[:3]):
        fx = 0.4 + i * 4.3
        fy, fw, fh = 1.0, 4.0, 5.8
        fc = fase_colors[i]
        add_rect(s5, fx, fy, fw, fh, C["bg2"])
        add_rect(s5, fx, fy, fw, 0.07, fc)
        add_rect(s5, fx + 0.15, fy + 0.15, fw - 0.3, 0.6, fc)
        add_text(s5, f.get("fase", ""), fx + 0.15, fy + 0.15, fw - 0.3, 0.6,
                 font_size=11, bold=True, color_hex=C["bg"], align="center")
        for j, a in enumerate((f.get("acciones") or [])[:4]):
            ay = fy + 0.95 + j * 0.75
            add_rect(s5, fx + 0.15, ay + 0.1, 0.3, 0.3, fc)
            add_text(s5, str(j + 1), fx + 0.15, ay + 0.1, 0.3, 0.3,
                     font_size=8, bold=True, color_hex=C["bg"], align="center")
            add_text(s5, a, fx + 0.55, ay, fw - 0.75, 0.6,
                     font_size=9, color_hex=C["white"])
        add_rect(s5, fx + 0.15, fy + fh - 1.05, fw - 0.3, 0.45, C["bg"])
        add_text(s5, f"💶 {f.get('inversion','')}", fx + 0.15, fy + fh - 1.05, fw - 0.3, 0.45,
                 font_size=10, bold=True, color_hex=fc, align="center")
        add_text(s5, f.get("riesgo_mitigado", ""), fx + 0.15, fy + fh - 0.55, fw - 0.3, 0.45,
                 font_size=8, color_hex=C["gray"], align="center")

    # ══ SLIDE 6 — CIERRE / LLAMADA A LA ACCIÓN ══
    s6 = add_slide()
    add_rect(s6, 0, 0, 13.3, 3.375, "1A0A0A")
    add_rect(s6, 0, 0, 0.5, 7.5, C["red"])
    add_text(s6, "⚠", 1, 0.3, 11.8, 1.0, font_size=40, align="center", color_hex=C["red"])
    add_text(s6, "EL RIESGO ES REAL. EL MOMENTO ES AHORA.",
             1, 1.2, 11.3, 1.0, font_size=24, bold=True, color_hex=C["white"], align="center")
    add_text(s6, data.get("llamada_accion", ""), 1.5, 2.4, 10.3, 0.9,
             font_size=16, color_hex=C["red"], italic=True, align="center")
    add_text(s6, str(score), 1, 3.7, 3, 2.0,
             font_size=72, bold=True, color_hex=C["red"], align="center")
    add_text(s6, "/100", 4.0, 4.8, 1.5, 0.7,
             font_size=20, color_hex=C["gray"])
    add_text(s6, "Score de Riesgo Global", 1, 5.65, 4, 0.4,
             font_size=10, color_hex=C["gray"], align="center")
    add_text(s6, "Próximos pasos:", 6, 3.7, 6.8, 0.45,
             font_size=13, bold=True, color_hex=C["white"])
    ctas = [
        "Aprobar presupuesto de remediación urgente",
        "Designar responsable de seguimiento (CISO / DPO)",
        "Iniciar acciones críticas en los próximos 30 días",
        "Programar revisión de progreso en 60 días",
    ]
    for i, c in enumerate(ctas):
        add_rect(s6, 6.1, 4.25 + i * 0.72, 0.3, 0.3, C["red"])
        add_text(s6, str(i + 1), 6.1, 4.25 + i * 0.72, 0.3, 0.3,
                 font_size=8, bold=True, color_hex=C["bg"], align="center")
        add_text(s6, c, 6.55, 4.2 + i * 0.72, 6.2, 0.45,
                 font_size=11, color_hex=C["white"])
    add_text(s6, f"Generado por AuditAI  ·  {data.get('fecha','')}",
             1, 7.05, 11.3, 0.35, font_size=8, color_hex=C["darkgray"], align="center")

    prs.save(output_path)


def generate_pptx_script(data: dict, output_path: str) -> str:
    """OBSOLETO — mantenido por compatibilidad. Usar generate_pptx_python."""
    import json as json_mod
    data_json = json_mod.dumps(data, ensure_ascii=False)

    return f"""
const pptxgen = require("pptxgenjs");

const data = {data_json};

// ── Paleta: Midnight Executive con acento rojo peligro ──
const C = {{
  bg:       "0D1117",
  bg2:      "161B22",
  bg3:      "21262D",
  red:      "FF4757",
  orange:   "FFA502",
  green:    "00E5A0",
  blue:     "5B9FFF",
  white:    "E8EAF0",
  gray:     "8892A4",
  darkgray: "4A5568",
  border:   "21262D",
}};

const CRIT_COLOR = {{ "Crítico": C.red, "Alto": C.orange, "Medio": "FFD700", "Bajo": C.green }};
const NIVEL_COLOR = {{ "CRÍTICO": C.red, "ALTO": C.orange, "MEDIO": "FFD700", "BAJO": C.green }};

let pres = new pptxgen();
pres.layout = "LAYOUT_WIDE"; // 13.3" x 7.5"
pres.title = data.titulo;
pres.author = "AuditAI";

const W = 13.3, H = 7.5;

// ════════════════════════════════════════════════
// SLIDE 1 — PORTADA
// ════════════════════════════════════════════════
{{
  let s = pres.addSlide();
  s.background = {{ color: C.bg }};

  // Bloque rojo izquierdo
  s.addShape(pres.shapes.RECTANGLE, {{ x:0, y:0, w:0.5, h:H, fill:{{ color: C.red }} }});

  // Nivel de riesgo badge grande
  const nivel = data.resumen_ejecutivo.nivel_riesgo;
  const nivelColor = NIVEL_COLOR[nivel] || C.red;
  s.addShape(pres.shapes.RECTANGLE, {{ x:0.9, y:1.0, w:3.2, h:0.65,
    fill:{{ color: nivelColor }}, rectRadius:0.05 }});
  s.addText("NIVEL DE RIESGO: " + nivel, {{ x:0.9, y:1.0, w:3.2, h:0.65,
    fontSize:14, bold:true, color:C.bg, align:"center", valign:"middle", margin:0 }});

  // Título principal
  s.addText(data.titulo, {{ x:0.9, y:1.85, w:9, h:1.3,
    fontSize:44, bold:true, color:C.white, fontFace:"Arial Black" }});
  s.addText(data.subtitulo, {{ x:0.9, y:3.1, w:9, h:0.6,
    fontSize:20, color:nivelColor, fontFace:"Arial" }});

  // Frase impacto
  s.addShape(pres.shapes.RECTANGLE, {{ x:0.9, y:3.9, w:8.5, h:0.8,
    fill:{{ color: C.bg2 }} }});
  s.addText("⚠  " + data.resumen_ejecutivo.frase_impacto, {{ x:0.9, y:3.9, w:8.5, h:0.8,
    fontSize:16, color:C.red, italic:true, align:"center", valign:"middle", margin:0 }});

  // Métricas rápidas en fila
  const mets = [
    {{ label: "Hallazgos Críticos", val: String(data.resumen_ejecutivo.hallazgos_criticos), color: C.red }},
    {{ label: "Hallazgos Altos",    val: String(data.resumen_ejecutivo.hallazgos_altos),    color: C.orange }},
    {{ label: "Hallazgos Medios",   val: String(data.resumen_ejecutivo.hallazgos_medios),   color: "FFD700" }},
    {{ label: "Score Global",       val: data.resumen_ejecutivo.score_global + "/100",      color: nivelColor }},
  ];
  mets.forEach((m, i) => {{
    const bx = 0.9 + i * 2.9;
    s.addShape(pres.shapes.RECTANGLE, {{ x:bx, y:4.9, w:2.6, h:1.3,
      fill:{{ color:C.bg2 }} }});
    s.addShape(pres.shapes.RECTANGLE, {{ x:bx, y:4.9, w:2.6, h:0.07,
      fill:{{ color:m.color }} }});
    s.addText(m.val, {{ x:bx, y:5.0, w:2.6, h:0.65,
      fontSize:32, bold:true, color:m.color, align:"center", valign:"middle", margin:0 }});
    s.addText(m.label, {{ x:bx, y:5.7, w:2.6, h:0.4,
      fontSize:10, color:C.gray, align:"center", valign:"middle", margin:0 }});
  }});

  // Fecha y empresa
  s.addText(data.empresa + "  ·  " + data.fecha, {{ x:0.9, y:H-0.5, w:10, h:0.35,
    fontSize:10, color:C.darkgray }});
}}

// ════════════════════════════════════════════════
// SLIDE 2 — TOP 5 RIESGOS CRÍTICOS (tabla visual)
// ════════════════════════════════════════════════
{{
  let s = pres.addSlide();
  s.background = {{ color: C.bg }};

  s.addShape(pres.shapes.RECTANGLE, {{ x:0, y:0, w:W, h:0.8, fill:{{ color:C.bg2 }} }});
  s.addText("TOP 5 RIESGOS CRÍTICOS", {{ x:0.5, y:0, w:W-1, h:0.8,
    fontSize:22, bold:true, color:C.white, valign:"middle", margin:0 }});
  s.addText("Impacto directo sobre el negocio si no se actúa", {{ x:0.5, y:0, w:W-1, h:0.8,
    fontSize:12, color:C.gray, align:"right", valign:"middle", margin:0 }});

  const risks = (data.top_riesgos || []).slice(0, 5);
  risks.forEach((r, i) => {{
    const ry = 0.95 + i * 1.28;
    const color = CRIT_COLOR[r.criticidad] || C.orange;

    // Fondo fila
    s.addShape(pres.shapes.RECTANGLE, {{ x:0.3, y:ry, w:W-0.6, h:1.15,
      fill:{{ color:C.bg2 }} }});
    // Borde izquierdo color criticidad
    s.addShape(pres.shapes.RECTANGLE, {{ x:0.3, y:ry, w:0.08, h:1.15,
      fill:{{ color:color }} }});
    // Badge ID
    s.addShape(pres.shapes.RECTANGLE, {{ x:0.55, y:ry+0.25, w:0.5, h:0.5,
      fill:{{ color:color }} }});
    s.addText(r.id, {{ x:0.55, y:ry+0.25, w:0.5, h:0.5,
      fontSize:12, bold:true, color:C.bg, align:"center", valign:"middle", margin:0 }});
    // Título riesgo
    s.addText(r.titulo, {{ x:1.2, y:ry+0.08, w:5.5, h:0.4,
      fontSize:14, bold:true, color:C.white, margin:0 }});
    // Descripción
    s.addText(r.descripcion, {{ x:1.2, y:ry+0.5, w:5.5, h:0.55,
      fontSize:10, color:C.gray, margin:0 }});
    // Impacto económico
    s.addText("💰 " + r.impacto_economico, {{ x:6.9, y:ry+0.1, w:3.5, h:0.4,
      fontSize:11, bold:true, color:color, margin:0 }});
    // Consecuencia
    s.addText("⚡ " + r.consecuencia, {{ x:6.9, y:ry+0.55, w:3.5, h:0.45,
      fontSize:10, color:C.gray, margin:0 }});
    // Badge criticidad
    s.addShape(pres.shapes.RECTANGLE, {{ x:10.6, y:ry+0.3, w:1.5, h:0.4,
      fill:{{ color:color }} }});
    s.addText(r.criticidad.toUpperCase(), {{ x:10.6, y:ry+0.3, w:1.5, h:0.4,
      fontSize:10, bold:true, color:C.bg, align:"center", valign:"middle", margin:0 }});
    // Normativas
    const normsText = (r.normativas || []).join(" · ");
    s.addText(normsText, {{ x:10.55, y:ry+0.75, w:1.6, h:0.3,
      fontSize:8, color:C.darkgray, align:"center", margin:0 }});
  }});
}}

// ════════════════════════════════════════════════
// SLIDE 3 — EXPOSICIÓN NORMATIVA
// ════════════════════════════════════════════════
{{
  let s = pres.addSlide();
  s.background = {{ color: C.bg }};

  s.addShape(pres.shapes.RECTANGLE, {{ x:0, y:0, w:W, h:0.8, fill:{{ color:C.bg2 }} }});
  s.addText("EXPOSICIÓN NORMATIVA Y SANCIONES", {{ x:0.5, y:0, w:W-1, h:0.8,
    fontSize:22, bold:true, color:C.white, valign:"middle", margin:0 }});

  const norms = (data.exposicion_normativa || []).slice(0, 5);
  norms.forEach((n, i) => {{
    const col = i % 3, row = Math.floor(i / 3);
    const cx = 0.4 + col * 4.3;
    const cy = 1.0 + row * 2.8;
    const cw = 3.9, ch = 2.5;
    const color = NIVEL_COLOR[n.nivel] || C.orange;

    s.addShape(pres.shapes.RECTANGLE, {{ x:cx, y:cy, w:cw, h:ch, fill:{{ color:C.bg2 }} }});
    s.addShape(pres.shapes.RECTANGLE, {{ x:cx, y:cy, w:cw, h:0.07, fill:{{ color:color }} }});

    // Nombre normativa
    s.addText(n.normativa, {{ x:cx+0.15, y:cy+0.12, w:cw-0.6, h:0.45,
      fontSize:16, bold:true, color:C.white, margin:0 }});
    // Badge nivel
    s.addShape(pres.shapes.RECTANGLE, {{ x:cx+cw-0.9, y:cy+0.12, w:0.75, h:0.35,
      fill:{{ color:color }} }});
    s.addText(n.nivel, {{ x:cx+cw-0.9, y:cy+0.12, w:0.75, h:0.35,
      fontSize:8, bold:true, color:C.bg, align:"center", valign:"middle", margin:0 }});
    // Multa
    s.addText("Multa máx: " + n.multa_max, {{ x:cx+0.15, y:cy+0.65, w:cw-0.3, h:0.35,
      fontSize:11, bold:true, color:color, margin:0 }});
    // Incumplimientos
    s.addText(String(n.incumplimientos) + " incumplimientos detectados", {{ x:cx+0.15, y:cy+1.0, w:cw-0.3, h:0.3,
      fontSize:10, color:C.gray, margin:0 }});
    // Descripción
    s.addText(n.descripcion, {{ x:cx+0.15, y:cy+1.35, w:cw-0.3, h:0.95,
      fontSize:9, color:C.gray, margin:0 }});
  }});
}}

// ════════════════════════════════════════════════
// SLIDE 4 — MÉTRICAS CLAVE (big numbers)
// ════════════════════════════════════════════════
{{
  let s = pres.addSlide();
  s.background = {{ color: C.bg }};

  s.addShape(pres.shapes.RECTANGLE, {{ x:0, y:0, w:W, h:0.8, fill:{{ color:C.bg2 }} }});
  s.addText("MÉTRICAS CLAVE DE EXPOSICIÓN", {{ x:0.5, y:0, w:W-1, h:0.8,
    fontSize:22, bold:true, color:C.white, valign:"middle", margin:0 }});

  const metrics = (data.metricas_clave || []).slice(0, 4);
  const mColors = [C.red, C.orange, C.blue, C.green];
  metrics.forEach((m, i) => {{
    const col = i % 2, row = Math.floor(i / 2);
    const mx = 0.5 + col * 6.4, my = 1.0 + row * 2.8;
    const mw = 6.0, mh = 2.5;
    s.addShape(pres.shapes.RECTANGLE, {{ x:mx, y:my, w:mw, h:mh, fill:{{ color:C.bg2 }} }});
    s.addShape(pres.shapes.RECTANGLE, {{ x:mx, y:my, w:mw, h:0.07, fill:{{ color:mColors[i] }} }});
    s.addText(m.valor, {{ x:mx, y:my+0.15, w:mw, h:1.2,
      fontSize:56, bold:true, color:mColors[i], align:"center", valign:"middle", margin:0 }});
    s.addText(m.metrica, {{ x:mx+0.2, y:my+1.4, w:mw-0.4, h:0.45,
      fontSize:16, bold:true, color:C.white, align:"center", margin:0 }});
    s.addText(m.contexto, {{ x:mx+0.2, y:my+1.88, w:mw-0.4, h:0.45,
      fontSize:11, color:C.gray, align:"center", margin:0 }});
  }});
}}

// ════════════════════════════════════════════════
// SLIDE 5 — PLAN DE ACCIÓN
// ════════════════════════════════════════════════
{{
  let s = pres.addSlide();
  s.background = {{ color: C.bg }};

  s.addShape(pres.shapes.RECTANGLE, {{ x:0, y:0, w:W, h:0.8, fill:{{ color:C.bg2 }} }});
  s.addText("PLAN DE ACCIÓN RECOMENDADO", {{ x:0.5, y:0, w:W-1, h:0.8,
    fontSize:22, bold:true, color:C.white, valign:"middle", margin:0 }});

  const faseColors = [C.red, C.orange, C.blue];
  const fases = (data.plan_accion || []).slice(0, 3);
  fases.forEach((f, i) => {{
    const fx = 0.4 + i * 4.3;
    const fy = 1.0, fw = 4.0, fh = 5.8;
    const color = faseColors[i];

    s.addShape(pres.shapes.RECTANGLE, {{ x:fx, y:fy, w:fw, h:fh, fill:{{ color:C.bg2 }} }});
    s.addShape(pres.shapes.RECTANGLE, {{ x:fx, y:fy, w:fw, h:0.07, fill:{{ color:color }} }});

    // Fase header
    s.addShape(pres.shapes.RECTANGLE, {{ x:fx+0.15, y:fy+0.15, w:fw-0.3, h:0.6,
      fill:{{ color:color }} }});
    s.addText(f.fase, {{ x:fx+0.15, y:fy+0.15, w:fw-0.3, h:0.6,
      fontSize:12, bold:true, color:C.bg, align:"center", valign:"middle", margin:0 }});

    // Acciones
    const acciones = (f.acciones || []).slice(0, 4);
    acciones.forEach((a, j) => {{
      const ay = fy + 0.95 + j * 0.75;
      s.addShape(pres.shapes.OVAL, {{ x:fx+0.15, y:ay+0.1, w:0.3, h:0.3, fill:{{ color:color }} }});
      s.addText(String(j+1), {{ x:fx+0.15, y:ay+0.1, w:0.3, h:0.3,
        fontSize:9, bold:true, color:C.bg, align:"center", valign:"middle", margin:0 }});
      s.addText(a, {{ x:fx+0.55, y:ay, w:fw-0.75, h:0.6,
        fontSize:10, color:C.white, margin:0 }});
    }});

    // Inversión
    s.addShape(pres.shapes.RECTANGLE, {{ x:fx+0.15, y:fy+fh-1.05, w:fw-0.3, h:0.45,
      fill:{{ color:"0D1117" }} }});
    s.addText("💶 " + f.inversion, {{ x:fx+0.15, y:fy+fh-1.05, w:fw-0.3, h:0.45,
      fontSize:11, bold:true, color:color, align:"center", valign:"middle", margin:0 }});
    s.addText(f.riesgo_mitigado, {{ x:fx+0.15, y:fy+fh-0.55, w:fw-0.3, h:0.45,
      fontSize:9, color:C.gray, align:"center", margin:0 }});
  }});
}}

// ════════════════════════════════════════════════
// SLIDE 6 — CIERRE / LLAMADA A LA ACCIÓN
// ════════════════════════════════════════════════
{{
  let s = pres.addSlide();
  s.background = {{ color: C.bg }};

  // Fondo rojo dramático mitad superior
  s.addShape(pres.shapes.RECTANGLE, {{ x:0, y:0, w:W, h:H*0.45, fill:{{ color:"1A0A0A" }} }});
  s.addShape(pres.shapes.RECTANGLE, {{ x:0, y:0, w:0.5, h:H, fill:{{ color:C.red }} }});

  s.addText("⚠", {{ x:1, y:0.3, w:W-1.5, h:1.0,
    fontSize:48, align:"center", color:C.red }});
  s.addText("EL RIESGO ES REAL. EL MOMENTO ES AHORA.", {{ x:1, y:1.2, w:W-1.5, h:1.0,
    fontSize:28, bold:true, color:C.white, align:"center", fontFace:"Arial Black", margin:0 }});

  s.addText(data.llamada_accion, {{ x:1.5, y:2.4, w:W-2.5, h:0.9,
    fontSize:18, color:C.red, italic:true, align:"center", margin:0 }});

  // Línea divisoria
  s.addShape(pres.shapes.LINE, {{ x:1, y:3.5, w:W-2, h:0,
    line:{{ color:C.border, width:1 }} }});

  // Score final grande
  s.addText(String(data.resumen_ejecutivo.score_global), {{ x:1, y:3.7, w:3, h:2.0,
    fontSize:96, bold:true, color:C.red, align:"center", valign:"middle", margin:0 }});
  s.addText("/100", {{ x:4.0, y:4.8, w:1.5, h:0.7,
    fontSize:24, color:C.gray, valign:"middle", margin:0 }});
  s.addText("Score de Riesgo Global", {{ x:1, y:5.65, w:4, h:0.4,
    fontSize:11, color:C.gray, align:"center", margin:0 }});

  // CTA texto
  s.addText("Próximos pasos:", {{ x:6, y:3.7, w:W-6.5, h:0.45,
    fontSize:14, bold:true, color:C.white, margin:0 }});

  const ctas = [
    "Aprobar presupuesto de remediación urgente",
    "Designar responsable de seguimiento (CISO / DPO)",
    "Iniciar acciones críticas en los próximos 30 días",
    "Programar revisión de progreso en 60 días"
  ];
  ctas.forEach((c, i) => {{
    s.addShape(pres.shapes.OVAL, {{ x:6.1, y:4.25+i*0.72, w:0.3, h:0.3, fill:{{ color:C.red }} }});
    s.addText(c, {{ x:6.55, y:4.2+i*0.72, w:W-7.1, h:0.45,
      fontSize:12, color:C.white, margin:0 }});
  }});

  s.addText("Generado por AuditAI  ·  " + data.fecha, {{ x:1, y:H-0.45, w:W-2, h:0.35,
    fontSize:9, color:C.darkgray, align:"center", margin:0 }});
}}

pres.writeFile({{ fileName: "{output_path}" }})
  .then(() => {{ console.log("OK:" + "{output_path}"); }})
  .catch(e => {{ console.error("ERROR:" + e); process.exit(1); }});
"""


def risk_data_to_html(data: dict) -> str:
    """Genera un HTML de preview del análisis de riesgos para mostrar en la pestaña."""
    nivel = data.get("resumen_ejecutivo", {}).get("nivel_riesgo", "ALTO")
    nivel_colors = {"CRÍTICO": "#ff4757", "ALTO": "#ffa502", "MEDIO": "#ffd700", "BAJO": "#00e5a0"}
    color = nivel_colors.get(nivel, "#ffa502")
    score = data.get("resumen_ejecutivo", {}).get("score_global", 0)
    frase = data.get("resumen_ejecutivo", {}).get("frase_impacto", "")

    top_riesgos_html = ""
    for r in (data.get("top_riesgos") or []):
        rc = nivel_colors.get(r.get("criticidad","").upper(), "#ffa502")
        norms = " · ".join(r.get("normativas") or [])
        top_riesgos_html += f"""
      <div class="risk-card">
        <div class="risk-header">
          <span class="risk-id">{r.get('id','')}</span>
          <span class="risk-title">{r.get('titulo','')}</span>
          <span class="badge" style="background:rgba(0,0,0,0.3);border:1px solid {rc};color:{rc}">{r.get('criticidad','')}</span>
        </div>
        <p class="risk-desc">{r.get('descripcion','')}</p>
        <div class="risk-meta">
          <span class="risk-impact">💰 {r.get('impacto_economico','')}</span>
          <span class="risk-conseq">⚡ {r.get('consecuencia','')}</span>
        </div>
        <div class="risk-norms">{norms}</div>
      </div>"""

    norms_html = ""
    norm_colors = {"CRÍTICO": "#ff4757","ALTO": "#ffa502","MEDIO": "#ffd700","BAJO": "#00e5a0"}
    for n in (data.get("exposicion_normativa") or []):
        nc = norm_colors.get(n.get("nivel","").upper(), "#ffa502")
        norms_html += f"""
      <div class="norm-card">
        <div class="norm-top" style="border-top:3px solid {nc}">
          <span class="norm-name">{n.get('normativa','')}</span>
          <span class="badge" style="background:{nc};color:#0d1117">{n.get('nivel','')}</span>
        </div>
        <div class="norm-multa" style="color:{nc}">Multa máx: {n.get('multa_max','')}</div>
        <div class="norm-inc">{n.get('incumplimientos',0)} incumplimientos</div>
        <p class="norm-desc">{n.get('descripcion','')}</p>
      </div>"""

    metricas_html = ""
    met_colors = ["#ff4757","#ffa502","#5b9fff","#00e5a0"]
    for i, m in enumerate((data.get("metricas_clave") or [])[:4]):
        mc = met_colors[i % len(met_colors)]
        metricas_html += f"""
      <div class="metric-card" style="border-top:3px solid {mc}">
        <div class="metric-val" style="color:{mc}">{m.get('valor','')}</div>
        <div class="metric-name">{m.get('metrica','')}</div>
        <div class="metric-ctx">{m.get('contexto','')}</div>
      </div>"""

    plan_html = ""
    fase_colors = ["#ff4757","#ffa502","#5b9fff"]
    for i, f in enumerate((data.get("plan_accion") or [])[:3]):
        fc = fase_colors[i % len(fase_colors)]
        acciones_li = "".join(f"<li>{a}</li>" for a in (f.get("acciones") or []))
        plan_html += f"""
      <div class="fase-card" style="border-top:3px solid {fc}">
        <div class="fase-header" style="background:{fc};color:#0d1117">{f.get('fase','')}</div>
        <ul class="fase-acciones">{acciones_li}</ul>
        <div class="fase-inv" style="color:{fc}">💶 {f.get('inversion','')}</div>
        <div class="fase-mit">{f.get('riesgo_mitigado','')}</div>
      </div>"""

    return f"""<!DOCTYPE html>
<html lang="es">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Análisis de Riesgos — AuditAI</title>
<link href="https://fonts.googleapis.com/css2?family=Syne:wght@700;800&family=Inter:wght@400;500&family=IBM+Plex+Mono:wght@400;500&display=swap" rel="stylesheet">
<style>
  :root {{ --bg:#0d1117;--bg2:#161b22;--bg3:#21262d;--red:#ff4757;--accent:{color};
           --text:#e8eaf0;--text2:#8892a4;--text3:#4a5568;--border:#21262d; }}
  *{{box-sizing:border-box;margin:0;padding:0}}
  html,body{{overflow-x:hidden}}
  body{{background:var(--bg);color:var(--text);font-family:'Inter',sans-serif;font-size:15px;line-height:1.7}}
  .page-wrap{{max-width:1100px;margin:0 auto;padding:50px 40px 100px}}
  .report-header{{background:var(--bg2);border:1px solid var(--border);border-radius:16px;padding:36px 40px;margin-bottom:48px;position:relative;overflow:hidden}}
  .report-header::before{{content:'⚠️';position:absolute;right:36px;top:50%;transform:translateY(-50%);font-size:80px;opacity:0.12}}
  .report-type{{font-family:'IBM Plex Mono',monospace;font-size:11px;color:{color};letter-spacing:2px;text-transform:uppercase;margin-bottom:10px}}
  h1{{font-family:'Syne',sans-serif;font-size:36px;font-weight:800;color:{color};margin-bottom:6px}}
  .report-meta{{display:flex;gap:20px;flex-wrap:wrap;margin-top:14px}}
  .report-meta span{{font-size:12px;color:var(--text3);font-family:'IBM Plex Mono',monospace}}
  .report-meta span b{{color:var(--text2)}}
  .nivel-badge{{display:inline-block;padding:4px 16px;border-radius:4px;font-size:13px;font-weight:700;letter-spacing:1px;background:{color};color:#0d1117;margin-bottom:12px}}
  .score-wrap{{display:flex;align-items:baseline;gap:6px;margin:10px 0}}
  .score-num{{font-size:56px;font-weight:800;color:{color};font-family:'Syne',sans-serif}}
  .score-label{{font-size:14px;color:var(--text2)}}
  .frase{{font-style:italic;color:{color};font-size:16px;margin-top:8px;padding:12px 18px;background:rgba(255,71,87,0.07);border-left:3px solid {color};border-radius:0 8px 8px 0}}
  h2{{font-family:'Syne',sans-serif;font-size:20px;font-weight:700;color:#fff;margin:48px 0 18px;padding-bottom:10px;border-bottom:2px solid var(--border)}}
  .badge{{display:inline-block;padding:2px 10px;border-radius:20px;font-size:11px;font-weight:600;white-space:nowrap}}
  /* Métricas */
  .metrics-grid{{display:grid;grid-template-columns:repeat(4,1fr);gap:16px;margin-bottom:16px}}
  .metric-card{{background:var(--bg2);border-radius:12px;padding:20px;}}
  .metric-val{{font-size:36px;font-weight:800;font-family:'Syne',sans-serif;margin-bottom:4px}}
  .metric-name{{font-size:13px;font-weight:600;color:var(--text);margin-bottom:4px}}
  .metric-ctx{{font-size:11px;color:var(--text2)}}
  /* Riesgos */
  .risk-card{{background:var(--bg2);border-radius:12px;padding:20px 24px;margin-bottom:12px;border-left:4px solid {color}}}
  .risk-header{{display:flex;align-items:center;gap:12px;margin-bottom:10px;flex-wrap:wrap}}
  .risk-id{{font-family:'IBM Plex Mono',monospace;font-size:12px;color:{color};font-weight:700;background:rgba(255,71,87,0.12);padding:2px 10px;border-radius:4px}}
  .risk-title{{font-weight:700;font-size:15px;color:#fff;flex:1}}
  .risk-desc{{color:var(--text2);font-size:13px;margin-bottom:10px}}
  .risk-meta{{display:flex;gap:20px;flex-wrap:wrap;margin-bottom:6px}}
  .risk-impact{{font-size:12px;font-weight:600;color:#ffa502}}
  .risk-conseq{{font-size:12px;color:var(--text2)}}
  .risk-norms{{font-size:11px;color:var(--text3);font-family:'IBM Plex Mono',monospace;margin-top:6px}}
  /* Normativas */
  .norms-grid{{display:grid;grid-template-columns:repeat(3,1fr);gap:16px}}
  .norm-card{{background:var(--bg2);border-radius:12px;padding:18px}}
  .norm-top{{display:flex;justify-content:space-between;align-items:center;margin-bottom:10px;padding-top:12px}}
  .norm-name{{font-weight:700;font-size:16px;color:#fff}}
  .norm-multa{{font-size:12px;font-weight:600;margin-bottom:4px}}
  .norm-inc{{font-size:11px;color:var(--text2);margin-bottom:8px}}
  .norm-desc{{font-size:12px;color:var(--text2)}}
  /* Plan */
  .plan-grid{{display:grid;grid-template-columns:repeat(3,1fr);gap:16px}}
  .fase-card{{background:var(--bg2);border-radius:12px;overflow:hidden}}
  .fase-header{{padding:10px 16px;font-size:13px;font-weight:700;text-align:center}}
  .fase-acciones{{padding:16px 16px 0;list-style:none}}
  .fase-acciones li{{font-size:12px;color:var(--text2);padding:6px 0;border-bottom:1px solid var(--border);padding-left:14px;position:relative}}
  .fase-acciones li::before{{content:'→';position:absolute;left:0;color:var(--text3)}}
  .fase-inv{{padding:12px 16px 4px;font-size:13px;font-weight:700}}
  .fase-mit{{padding:0 16px 16px;font-size:11px;color:var(--text2)}}
  .cta-box{{background:linear-gradient(135deg,#1a0a0a,#0d1117);border:1px solid {color};border-radius:16px;padding:40px;text-align:center;margin-top:48px}}
  .cta-text{{font-family:'Syne',sans-serif;font-size:22px;font-weight:800;color:{color};margin-bottom:8px}}
  @media(max-width:800px){{.metrics-grid,.norms-grid,.plan-grid{{grid-template-columns:1fr}}}}
</style>
</head>
<body>
<div class="page-wrap">
  <div class="report-header">
    <div class="report-type">AuditAI — Análisis de Riesgos</div>
    <div class="nivel-badge">NIVEL: {nivel}</div>
    <h1>{data.get('titulo','Análisis de Riesgos')}</h1>
    <p style="color:var(--text2);margin:4px 0 10px">{data.get('subtitulo','')}</p>
    <div class="score-wrap">
      <div class="score-num">{score}</div>
      <div class="score-label">/100 — Score de Riesgo Global</div>
    </div>
    <div class="frase">{frase}</div>
    <div class="report-meta">
      <span><b>Sistema:</b> AuditAI Multi-Agente (CrewAI + Ollama)</span>
      <span><b>Fecha:</b> {data.get('fecha','')}</span>
    </div>
  </div>

  <h2>📊 Métricas Clave</h2>
  <div class="metrics-grid">{metricas_html}</div>

  <h2>🔴 Top 5 Riesgos Críticos</h2>
  {top_riesgos_html}

  <h2>⚖️ Exposición Normativa</h2>
  <div class="norms-grid">{norms_html}</div>

  <h2>🗓️ Plan de Acción</h2>
  <div class="plan-grid">{plan_html}</div>

  <div class="cta-box">
    <div class="cta-text">⚠ {data.get('llamada_accion','El riesgo es real. El momento es ahora.')}</div>
    <p style="color:var(--text2);margin-top:10px;font-size:14px">Descarga la presentación PowerPoint para presentar ante la dirección.</p>
  </div>
</div>
</body>
</html>"""


@app.post("/api/riesgos/generar-ppt")
async def generar_ppt_riesgos():
    """Lanza el agente de análisis de riesgos para generar el PPT."""
    md_clas = DB_DIR / "informe_clasificacion.md"
    md_ciso = DB_DIR / "informe_ciso.md"
    if not md_clas.exists() and not md_ciso.exists():
        raise HTTPException(400, "No hay informes generados. Ejecuta primero el análisis principal.")

    job_id = str(uuid.uuid4())
    ppt_jobs[job_id] = {
        "estado": "starting",
        "mensaje": "Iniciando agente de análisis de riesgos…",
        "error": None,
        "fecha": None,
        "data": None,
    }
    threading.Thread(target=run_risk_ppt, args=(job_id,), daemon=True).start()
    return {"job_id": job_id}


@app.get("/api/riesgos/estado/{job_id}")
async def estado_ppt_job(job_id: str):
    if job_id not in ppt_jobs:
        raise HTTPException(404, "Job no encontrado")
    job = ppt_jobs[job_id]
    return {
        "estado": job["estado"],
        "mensaje": job.get("mensaje", ""),
        "error": job.get("error"),
        "fecha": job.get("fecha"),
    }


@app.get("/api/riesgos/html", response_class=HTMLResponse)
async def riesgos_html():
    path = DB_DIR / "informe_riesgos.html"
    if not path.exists():
        raise HTTPException(404, "No hay presentación de riesgos generada todavía")
    return HTMLResponse(path.read_text(encoding="utf-8"))


@app.get("/api/riesgos/descargar")
async def descargar_ppt():
    path = DB_DIR / "informe_riesgos.pptx"
    if not path.exists():
        raise HTTPException(404, "No hay PPTX generado todavía")
    fecha = datetime.now().strftime("%Y%m%d")
    return FileResponse(path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        filename=f"analisis_riesgos_{fecha}.pptx")
